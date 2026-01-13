from flask import Flask, request, jsonify, send_from_directory, abort
from pptx import Presentation
import docx
import joblib
import os
import time
from functools import wraps
from collections import defaultdict

app = Flask(__name__, static_folder=".", static_url_path="")

# Завантажуємо модель
model = joblib.load("text_model.pkl")
vectorizer = joblib.load("text_vectorizer.pkl")

# Захист: rate limit по IP
VISITS = defaultdict(list)
MAX_CALLS = 5          # максимум 5 перевірок
PER_SECONDS = 60       # за 1 хвилину

MAX_TEXT_SIZE = 10000  # обмеження на розмір тексту

# ---------------- Helper functions ----------------

def rate_limit(func):
    @wraps(func)
    def wrapper(*args, **kwargs):
        ip = request.remote_addr
        now = time.time()
        VISITS[ip] = [t for t in VISITS[ip] if now - t < PER_SECONDS]
        if len(VISITS[ip]) >= MAX_CALLS:
            return jsonify({"error":"Занадто багато запитів. Зачекайте."}), 429
        VISITS[ip].append(now)
        return func(*args, **kwargs)
    return wrapper

def extract_text(file):
    filename = file.filename.lower()
    if filename.endswith(".pptx"):
        prs = Presentation(file)
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    elif filename.endswith(".docx"):
        doc = docx.Document(file)
        return "\n".join([p.text for p in doc.paragraphs])
    elif filename.endswith(".txt"):
        text = file.read()
        if isinstance(text, bytes):
            text = text.decode("utf-8", errors="ignore")
        return text
    return ""

# ---------------- Routes ----------------

@app.route("/")
def home():
    return send_from_directory(".", "index.html")

@app.route("/predict", methods=["POST"])
@rate_limit
def predict():
    text = ""

    if "file" in request.files:
        text = extract_text(request.files["file"])
    elif request.json and "text" in request.json:
        text = request.json["text"]

    if not text or len(text.strip()) < 10:
        return jsonify({"error":"Текст занадто короткий або порожній"}), 400

    if len(text) > MAX_TEXT_SIZE:
        return jsonify({"error":"Текст занадто великий"}), 400

    # Обробка через модель
    X = vectorizer.transform([text])
    proba = model.predict_proba(X)[0]
    ai_prob = float(proba[1])

    return jsonify({
        "ai_percent": round(ai_prob * 100,1),
        "verdict": "AI Generated" if ai_prob >= 0.5 else "Human Written"
    })

# ---------------- Run ----------------

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port)
